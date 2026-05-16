from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

_SERVER_ROOT = Path(__file__).resolve().parents[2]
_TEMPLATES_DIR = _SERVER_ROOT / "assets" / "templates"


@dataclass(frozen=True)
class ThemeSpec:
    id: str
    label: str
    file: str
    background: RGBColor
    accent: RGBColor
    title_color: RGBColor
    head_color: RGBColor
    detail_color: RGBColor
    strip_sample_slides: bool = False
    use_accent_bar: bool = True
    preserve_layout: bool = False


THEMES: dict[str, ThemeSpec] = {
    "brand": ThemeSpec(
        id="brand",
        label="Untitled Presentation",
        file="untitled-presentation.pptx",
        background=RGBColor(0xFF, 0xFF, 0xFF),
        accent=RGBColor(0x25, 0x63, 0xEB),
        title_color=RGBColor(0x1A, 0x36, 0x5D),
        head_color=RGBColor(0x1E, 0x3A, 0x5F),
        detail_color=RGBColor(0x47, 0x55, 0x69),
        strip_sample_slides=True,
        use_accent_bar=False,
        preserve_layout=True,
    ),
    "classroom": ThemeSpec(
        id="classroom",
        label="Classroom Blue",
        file="classroom.pptx",
        background=RGBColor(0xF5, 0xF8, 0xFF),
        accent=RGBColor(0x25, 0x63, 0xEB),
        title_color=RGBColor(0x1A, 0x36, 0x5D),
        head_color=RGBColor(0x1E, 0x3A, 0x5F),
        detail_color=RGBColor(0x47, 0x55, 0x69),
    ),
    "ocean": ThemeSpec(
        id="ocean",
        label="Ocean Teal",
        file="ocean.pptx",
        background=RGBColor(0xE8, 0xF6, 0xFA),
        accent=RGBColor(0x0D, 0x94, 0x88),
        title_color=RGBColor(0x0F, 0x4C, 0x5C),
        head_color=RGBColor(0x0F, 0x4C, 0x5C),
        detail_color=RGBColor(0x33, 0x5C, 0x67),
    ),
    "forest": ThemeSpec(
        id="forest",
        label="Forest Green",
        file="forest.pptx",
        background=RGBColor(0xEC, 0xF7, 0xF0),
        accent=RGBColor(0x16, 0x8A, 0x45),
        title_color=RGBColor(0x14, 0x53, 0x2D),
        head_color=RGBColor(0x14, 0x53, 0x2D),
        detail_color=RGBColor(0x3D, 0x5C, 0x45),
    ),
}

DEFAULT_THEME = "brand"


def list_themes() -> list[dict[str, str]]:
    order = ("brand", "classroom", "ocean", "forest")
    out: list[dict[str, str]] = []
    for key in order:
        if key in THEMES:
            t = THEMES[key]
            out.append({"id": t.id, "label": t.label})
    for t in THEMES.values():
        if t.id not in order:
            out.append({"id": t.id, "label": t.label})
    return out


def resolve_theme(theme: str | None) -> ThemeSpec:
    key = (theme or DEFAULT_THEME).strip().lower()
    if key in THEMES:
        return THEMES[key]
    custom = _TEMPLATES_DIR / f"{key}.pptx"
    if custom.is_file():
        base = THEMES[DEFAULT_THEME]
        return ThemeSpec(
            id=key,
            label=key.replace("-", " ").title(),
            file=f"{key}.pptx",
            background=base.background,
            accent=base.accent,
            title_color=base.title_color,
            head_color=base.head_color,
            detail_color=base.detail_color,
            strip_sample_slides=True,
            use_accent_bar=False,
            preserve_layout=True,
        )
    return THEMES[DEFAULT_THEME]


def templates_dir() -> Path:
    return _TEMPLATES_DIR


def template_path(spec: ThemeSpec) -> Path:
    return _TEMPLATES_DIR / spec.file


def _write_blank_template(path: Path, spec: ThemeSpec) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    fill = prs.slide_master.background.fill
    fill.solid()
    fill.fore_color.rgb = spec.background
    prs.save(str(path))


def _clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def ensure_template(spec: ThemeSpec) -> Path:
    path = template_path(spec)
    if path.is_file():
        return path
    if spec.id not in THEMES:
        raise FileNotFoundError(f"Template not found: {path}")
    _write_blank_template(path, spec)
    return path


def open_themed_presentation(theme: str | None) -> tuple[Presentation, ThemeSpec]:
    spec = resolve_theme(theme)
    path = ensure_template(spec)
    prs = Presentation(str(path))
    if spec.strip_sample_slides and len(prs.slides) > 0:
        _clear_slides(prs)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, spec


def layout_for(prs: Presentation, kind: str):
    want = {
        "title": ("title slide", "title only"),
        "content": ("title and content", "title & content", "content"),
    }.get(kind, ())
    for layout in prs.slide_layouts:
        name = layout.name.lower()
        if any(w in name for w in want):
            return layout
    idx = 0 if kind == "title" else 1
    return prs.slide_layouts[min(idx, len(prs.slide_layouts) - 1)]
