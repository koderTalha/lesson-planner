import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.api.slide_plan_models import SlidePlan
from app.services.pptx_themes import ThemeSpec, open_themed_presentation, layout_for

_FONT = "Calibri"
_MAX_DECK_TITLE = 72
_MAX_DECK_SUB = 100
_MAX_SLIDE_TITLE = 64
_MAX_GROUPS = 4
_MAX_HEAD_CHARS = 55
_MAX_DETAIL_CHARS = 130
_ACCENT_BAR_H = Inches(0.14)


def _truncate(text: str, limit: int) -> str:
    t = re.sub(r"\s+", " ", (text or "").strip())
    if len(t) <= limit:
        return t
    cut = t[: limit - 1].rsplit(" ", 1)[0].strip()
    return (cut or t[: limit - 1]).rstrip(".,;:") + "…"


def _bullet_groups(bullets: list) -> list[tuple[str, str | None]]:
    groups: list[tuple[str, str | None]] = []
    for b in bullets[:_MAX_GROUPS]:
        head = _truncate(str(getattr(b, "text", b)), _MAX_HEAD_CHARS)
        detail_raw = getattr(b, "detail", None)
        detail = _truncate(str(detail_raw), _MAX_DETAIL_CHARS) if detail_raw else None
        if not head:
            continue
        if not detail and ": " in head:
            parts = head.split(": ", 1)
            head, detail = _truncate(parts[0], _MAX_HEAD_CHARS), _truncate(
                parts[1], _MAX_DETAIL_CHARS
            )
        groups.append((head, detail))
    return groups


def _set_frame_margins(tf, *, left=0.12, right=0.12, top=0.08, bottom=0.08) -> None:
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)


def _style_paragraph(
    p,
    *,
    spec: ThemeSpec,
    size: int,
    bold: bool = False,
    color: RGBColor | None = None,
    space_after: int = 4,
    level: int = 0,
) -> None:
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    p.line_spacing = 1.12
    if not spec.preserve_layout:
        font = p.font
        font.name = _FONT
        font.size = Pt(size)
        font.bold = bold
        font.color.rgb = color or spec.head_color


def _fill_title(tf, text: str, *, spec: ThemeSpec, size: int) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if not spec.preserve_layout:
        _set_frame_margins(tf, top=0.05, bottom=0.05)
    p = tf.paragraphs[0]
    p.text = text
    _style_paragraph(
        p, spec=spec, size=size, bold=True, color=spec.title_color, space_after=0
    )


def _fill_bullet_groups(
    tf, groups: list[tuple[str, str | None]], *, spec: ThemeSpec
) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if not spec.preserve_layout:
        _set_frame_margins(tf)
    if not groups:
        p = tf.paragraphs[0]
        p.text = "—"
        _style_paragraph(p, spec=spec, size=18)
        return
    head_size = 19 if len(groups) <= 3 else 17
    detail_size = 16 if len(groups) <= 3 else 15
    first = True
    for head, detail in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = head
        _style_paragraph(
            p,
            spec=spec,
            size=head_size,
            bold=True,
            color=spec.head_color,
            space_after=2 if detail else 8,
            level=0,
        )
        if detail:
            d = tf.add_paragraph()
            d.text = detail
            _style_paragraph(
                d,
                spec=spec,
                size=detail_size,
                bold=False,
                color=spec.detail_color,
                space_after=10,
                level=1,
            )


def _accent_bar(slide, spec: ThemeSpec) -> None:
    if not spec.use_accent_bar:
        return
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        _ACCENT_BAR_H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = spec.accent
    bar.line.fill.background()


def _add_title_slide(prs: Presentation, plan: SlidePlan, spec: ThemeSpec) -> None:
    slide = prs.slides.add_slide(layout_for(prs, "title"))
    _accent_bar(slide, spec)
    title = _truncate(plan.deck_title, _MAX_DECK_TITLE)
    _fill_title(slide.shapes.title.text_frame, title, spec=spec, size=40)
    sub = _truncate(plan.deck_subtitle, _MAX_DECK_SUB)
    if sub and len(slide.placeholders) > 1:
        _fill_title(slide.placeholders[1].text_frame, sub, spec=spec, size=22)


def _add_content_slide(
    prs: Presentation, title: str, bullets: list, notes: str, spec: ThemeSpec
) -> None:
    slide = prs.slides.add_slide(layout_for(prs, "content"))
    _accent_bar(slide, spec)
    groups = _bullet_groups(bullets)
    _fill_title(
        slide.shapes.title.text_frame,
        _truncate(title, _MAX_SLIDE_TITLE),
        spec=spec,
        size=30,
    )
    body_ph = slide.placeholders[1]
    if not spec.preserve_layout:
        body_ph.top = Inches(1.2)
        body_ph.height = Inches(5.85)
        body_ph.width = Inches(9.0)
        body_ph.left = Inches(0.5)
    _fill_bullet_groups(body_ph.text_frame, groups, spec=spec)
    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()[:1200]


def build_pptx_from_plan(plan: SlidePlan, *, theme: str | None = None) -> bytes:
    prs, spec = open_themed_presentation(theme)
    _add_title_slide(prs, plan, spec)
    for slide in plan.slides:
        _add_content_slide(prs, slide.title, slide.bullets, slide.speaker_notes, spec)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
